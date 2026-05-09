"""
slide_insight_adder.py
======================
Utilities to *read* category labels from a PowerPoint deck, *match* them to
model-generated insights, and *write* those insights back into slide subheaders.
This module is intentionally thin and presentation-focused. It does *not*
generate insights; it only maps and applies text to slides.

Key ideas
---------
- We parse slide placeholders to discover category text (e.g., "Total Outerwear").
- We map those strings to a canonical list from your analysis DataFrame
  (case-insensitive, whitespace-insensitive, partial-match tolerant).
- We write either:
    1) Category-level bullets ("Category: insight") — optional legacy flow, or
    2) A single meta insight per slide (preferred), or
    3) A special *total* slide subheader.

Conventions
-----------
- Placeholder index **14** is used for subheaders in your template.
  If your template changes, update that index in your calling code or
  pass a different `placeholder_idx`.

- Slide keys are strings like `"Slide 8"`. Internally, we convert them to
  zero-based indices when addressing `prs.slides`.

- All write operations return a *status dict* to help with diagnostics.

Typical flow
------------
1) Use the analysis code to produce a DataFrame of insights per category.
2) Use `map_categories_from_placeholders()` to discover category labels per slide.
3) Use `filter_slide_mapping()` to keep only categories present in your analysis.
4) Build a category->insight lookup with `build_insight_lookup()`.
5) (Option A) `populate_subheadings_from_mapping()` to write per-category bullets,
   or (Option B) prepare a `meta_df` with one row per slide and call
   `apply_meta_insights()` to write a single meta subheader.
6) For the total slide, call `apply_total_subheader_to_slide()` with a
   pre-computed topline subheader string.

Dependencies
------------
- python-pptx (Presentation, shapes/placeholders API)
- pandas (DataFrame handling)

This file aims to be safe-by-default:
- If we can't find the exact placeholder by idx, we fall back to a reasonable
  text box target (subtitle, other text-capable placeholder, or a new textbox).
- All public entry points return useful status dictionaries instead of raising,
  unless inputs are structurally invalid.
"""
# =============================================================================
#Imports
# =============================================================================

from typing import Dict, List, Mapping, Optional, Tuple
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.util import Pt, Cm

# =============================================================================
# Category mapping / filtering
# =============================================================================

def filter_slide_mapping(slide_cats, collapsed):
    """
    Filter slide->text mapping down to *recognized* categories from `collapsed`.
    Matching is case- and whitespace-insensitive and allows partial matches:
    e.g., "total outerwear" will match "Total Outerwear".
    Parameters
    ----------
    slide_cats : dict[str, list[str]]
        From `map_categories_from_placeholders()`, e.g.:
        {'Slide 4': ['Total Sleepwear', 'Total Outerwear'], ...}
    collapsed : pandas.DataFrame
        Must contain a 'category' column with canonical category names.
    Returns
    -------
    dict[str, list[str]]
        For each slide key ('Slide N'), a list of matched, deduplicated category
        names exactly as they appear in `collapsed['category']`.
    """
    if "category" not in collapsed.columns:
        raise ValueError("Expected 'category' column in collapsed DataFrame")
    # Normalize valid category names: keep original, index by lowercased key.
    valid_original = [str(c).strip() for c in collapsed["category"].dropna().unique()]
    valid_norm_map = {v.lower(): v for v in valid_original}

    filtered = {}

    for slide, texts in slide_cats.items():
        matched = []
        seen = set()
        for text in texts:
            tclean = str(text).strip().lower()
            # Check each known category for partial match
            for vnorm, original in valid_norm_map.items():
                if vnorm in tclean and vnorm not in seen:
                    matched.append(original)
                    seen.add(vnorm)

        if matched:
            filtered[slide] = matched

    return filtered

# =============================================================================
# PowerPoint placeholder parsing
# =============================================================================
def map_categories_from_placeholders(prs):
    """
    Scan slides to collect likely category labels from BODY placeholders.
    Skips TITLE/SUBTITLE and the known subheading placeholder (idx=14).
    Returns the *first non-empty line* of any text-capable placeholder as a
    candidate category label.
    Returns
    dict[str, list[str]]
        Example:
        {
            'Slide 4': ['Total Bodysuits/Leotards', 'Total Childrens/Infants Sets', 'Total Sleepwear'],
            'Slide 5': ['Total Socks/Hosiery', 'Total Special Infantswear', 'Total Sportswear'],
            ...
        }
    """
    category_slide_map = {}

    for slide_idx, slide in enumerate(prs.slides, start=1):
        texts = []

        for shape in slide.shapes:
            try:
                if not getattr(shape, "is_placeholder", False):
                    continue
                if not getattr(shape, "has_text_frame", False):
                    continue

                ph_type = shape.placeholder_format.type
                ph_idx = getattr(shape.placeholder_format, "idx", -1)

                # Skip title/subtitle or the designated subheading placeholder
                if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.SUBTITLE):
                    continue
                if ph_idx == 14:  # reserved subheader placeholder in your template
                    continue

                raw_text = (shape.text_frame.text or "").strip()
                if not raw_text:
                    continue

                # Normalize to first non-empty line to reduce noise
                first_line = next((ln.strip() for ln in raw_text.splitlines() if ln.strip()), "")
                if first_line:
                    texts.append(first_line)

            except Exception:
                # Be resilient to odd shapes or theme quirks
                continue
        if texts:
            category_slide_map[f"Slide {slide_idx}"] = texts

    return category_slide_map

# =============================================================================
# Lookups / helpers
# =============================================================================

def build_insight_lookup(
    insights_df: pd.DataFrame,
    cat_col: str = "category",
    insight_col: str = "insight",
) -> Dict[str, str]:
    """
    Build a {category: insight} mapping from a DataFrame.
    - Drops NaNs
    - Keeps the *first* occurrence per category

    Raises ValueError if required columns are missing.
    """
    if cat_col not in insights_df.columns or insight_col not in insights_df.columns:
        raise ValueError(f"Expected columns '{cat_col}' and '{insight_col}' in insights_df.")
    cleaned = (
        insights_df[[cat_col, insight_col]]
        .dropna(subset=[cat_col, insight_col])
        .drop_duplicates(subset=[cat_col])
    )
    return dict(zip(cleaned[cat_col].astype(str), cleaned[insight_col].astype(str)))


def parse_slide_number_key(slide_key: str) -> Optional[int]:
    """
    Parse a key like 'Slide 8' to a zero-based slide index (7).
    Returns None if parsing fails.
    """
    try:
        num = int(str(slide_key).strip().split()[-1])
        return max(0, num - 1)
    except Exception:
        return None


def get_text_placeholder_by_idx(slide, idx: int):
    """
    Return a text-capable placeholder by *placeholder_format.idx*.
    Note: This is *not* the same as slide.placeholders[i] (positional).
    We walk shapes to find a placeholder whose `placeholder_format.idx == idx`.
    """
    for sh in slide.shapes:
        try:
            if getattr(sh, "is_placeholder", False) and hasattr(sh, "placeholder_format"):
                if sh.placeholder_format.idx == idx and hasattr(sh, "text_frame"):
                    return sh
        except Exception:
            pass
    return None


def ensure_subheading_shape(slide, preferred_placeholder_idx: int = 14):
    """
    Find a suitable text target for subheader/bullets on a slide.
    Preference order:
      1) Placeholder with idx == preferred_placeholder_idx
      2) SUBTITLE placeholder
      3) Other text-capable placeholders (non-TITLE) with short/empty text
      4) Any text box with short/empty text
      5) Create a new textbox at a sensible location
    Returns
    -------
    shape
        A shape with a `.text_frame` suitable for writing.
    """
    # 1) Known subheading placeholder
    sh = get_text_placeholder_by_idx(slide, preferred_placeholder_idx)
    if sh is not None:
        return sh

    # 2) Subtitle placeholder
    for s in slide.shapes:
        try:
            if s.is_placeholder and s.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE and hasattr(s, "text_frame"):
                return s
        except Exception:
            pass

    # 3) Other text-capable placeholders (avoid TITLE/CENTER_TITLE)
    for s in slide.shapes:
        try:
            if s.is_placeholder and hasattr(s, "text_frame"):
                if s.placeholder_format.type not in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    txt = (s.text_frame.text or "").strip()
                    if len(txt) <= 120:
                        return s
        except Exception:
            pass

    # 4) Any short/empty text box
    for s in slide.shapes:
        if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and hasattr(s, "text_frame"):
            txt = (s.text_frame.text or "").strip()
            if len(txt) <= 120:
                return s

    # 5) Create a new text box if all else fails
    left, top, width, height = Cm(2.0), Cm(4.0), Cm(24.0), Cm(6.0)
    return slide.shapes.add_textbox(left, top, width, height)


# =============================================================================
# Write meta insights (one line per slide) — PREFERRED
# =============================================================================

def apply_meta_insights(prs, meta_df, placeholder_idx=14):
    """
    Write a single synthesized meta insight into each slide's subheader.
    Expects a DataFrame with columns:
      - 'slide_id'     (e.g., 'Slide 8')
      - 'meta_insight' (plain text)
    Returns
    -------
    dict
        { 'Slide N': {'written': 0|1, 'missing': [...], 'target_idx': int|None, 'notes': str}, ... }
    """
    # logs the workings of the function to the console
    status = {}

    if meta_df is None or meta_df.empty:
        return status  # empty dict keeps callers safe

    for _, row in meta_df.iterrows():
        slide_id = row.get("slide_id")
        meta_text = (row.get("meta_insight") or "").strip()
        slide_key = str(slide_id)

        st = {"written": 0, "missing": [], "target_idx": placeholder_idx, "notes": ""}
        status[slide_key] = st

        # Basic validation
        if not slide_key or not meta_text:
            st["notes"] = "Missing slide_id or meta_insight."
            continue

        # Parse 'Slide 8' -> zero-based index
        try:
            slide_index = int(str(slide_key).split()[-1]) - 1
        except Exception:
            st["notes"] = f"Could not parse slide index from '{slide_key}'."
            continue

        if slide_index < 0 or slide_index >= len(prs.slides):
            st["notes"] = f"Slide index out of range for '{slide_key}'."
            continue

        slide = prs.slides[slide_index]
        try:
            # Prefer exact placeholder idx; otherwise pick a sensible text target
            ph = get_text_placeholder_by_idx(slide, placeholder_idx)
            if ph is None:
                ph = ensure_subheading_shape(slide, preferred_placeholder_idx=placeholder_idx)

            # Write the single-line meta insight
            if hasattr(ph, "text_frame"):
                tf = ph.text_frame
                # Keep first paragraph; drop extras
                while len(tf.paragraphs) > 1:
                    tf._element.remove(tf.paragraphs[-1]._p)
                tf.paragraphs[0].text = meta_text
                for run in tf.paragraphs[0].runs:
                    run.font.name = "Roboto Condensed"

                st["written"] = 1
                st["notes"] = "Meta insight applied"
                # Record actual target idx if it's a placeholder
                try:
                    st["target_idx"] = ph.placeholder_format.idx if ph.is_placeholder else None
                except Exception:
                    pass
            else:
                st["notes"] = "Target shape is not text-capable."
        except Exception as e:
            st["notes"] = f"Error applying meta insight: {e}"

    return status

# =============================================================================
# Total slide subheader (single line; uses same logic as meta)
# =============================================================================

def apply_total_subheader_to_slide(
    prs,
    total_slide,
    subheader_text: str,
    placeholder_idx: int = 14,
) -> dict:
    """
    Write the generated *total* subheader to the provided slide.
    Internally converts the slide to a 'Slide N' key and reuses `apply_meta_insights`.
    Returns
    -------
    dict
        Status dict with the same shape as `apply_meta_insights`.
    """
    # Convert slide object to 'Slide N' (1-based) for consistency with the API
    slide_num = prs.slides.index(total_slide) + 1
    slide_id = f"Slide {slide_num}"
    meta_df = pd.DataFrame([{"slide_id": slide_id, "meta_insight": subheader_text}])
    return apply_meta_insights(prs=prs, meta_df=meta_df, placeholder_idx=placeholder_idx)

# =============================================================================
# (Prev versio Keep here for now) Category bullets per slide — legacy/alternative flow
# =============================================================================
def replace_with_bullets(shape, lines: List[str], font_size_pt: Optional[float] = None):
    """
    Write each line as a bullet. Uses a literal bullet char for theme-agnostic output.
    Keeps the first paragraph to retain base formatting from the template.
    """
    tf = shape.text_frame
    # Remove all but the first paragraph
    while len(tf.paragraphs) > 1:
        tf._element.remove(tf.paragraphs[-1]._p)
    # Clear first paragraph
    p0 = tf.paragraphs[0]
    p0.text = ""
    def _set_font_size(paragraph):
        if font_size_pt is None:
            return
        for run in paragraph.runs:
            run.font.size = Pt(font_size_pt)
    # Write bullets
    for i, line in enumerate(lines):
        p = p0 if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        _set_font_size(p)


def populate_subheadings_from_mapping(
    prs: Presentation,
    slide_to_categories: Mapping[str, List[str]],
    insight_by_cat: Mapping[str, str],
    *,
    default_placeholder_idx: int = 14,
    per_slide_placeholder_idx: Optional[Mapping[str, int]] = None,
    font_size_pt: Optional[float] = None,
) -> Dict[str, dict]:
    """
    Write 'Category: insight' bullets into the subheading area of each slide.
    This is an alternative to `apply_meta_insights` and results in a bulleted list.
    Parameters
    ----------
    prs : Presentation
        python-pptx Presentation object.
    slide_to_categories : Mapping[str, List[str]]
        Mapping like {'Slide 4': ['Cat A', 'Cat B'], ...}
    insight_by_cat : Mapping[str, str]
        {'Cat A': 'Insight...', 'Cat B': 'Insight...'}
    default_placeholder_idx : int
        Preferred placeholder idx used when no per-slide override is supplied.
    per_slide_placeholder_idx : Mapping[str, int] | None
        Optional map like {'Slide 4': 14, 'Slide 5': 12, ...}
    font_size_pt : float | None
        Override font size for the bullets (None = keep theme default).

    Returns
    -------
    Dict[str, dict]
        Per-slide status, including counts written and any missing categories.
    """
    status: Dict[str, dict] = {}
    per_slide_placeholder_idx = per_slide_placeholder_idx or {}
    # Normalize lookup: case-insensitive category keys
    norm_insight = {str(k).strip().lower(): str(v) for k, v in insight_by_cat.items()}

    for slide_key, cats in slide_to_categories.items():
        per = {"written": 0, "missing": [], "target_idx": None, "notes": ""}
        idx0 = parse_slide_number_key(slide_key)
        if idx0 is None or idx0 < 0 or idx0 >= len(prs.slides):
            per["notes"] = "invalid slide number key"
            status[slide_key] = per
            continue

        slide = prs.slides[idx0]
        preferred_idx = per_slide_placeholder_idx.get(slide_key, default_placeholder_idx)
        shape = ensure_subheading_shape(slide, preferred_placeholder_idx=preferred_idx)

        # Record actual target idx (if placeholder) for debugging
        try:
            per["target_idx"] = shape.placeholder_format.idx if getattr(shape, "is_placeholder", False) else None
        except Exception:
            per["target_idx"] = None

        # De-dupe categories while preserving order
        seen = set()
        ordered = []
        for c in cats:
            key = str(c).strip()
            low = key.lower()
            if low not in seen:
                seen.add(low)
                ordered.append(key)
        # Compose "Category: insight" bullets
        lines = []
        for cat in ordered:
            insight = norm_insight.get(cat.strip().lower())
            if insight:
                lines.append(f"{cat}: {insight}")
            else:
                per["missing"].append(cat)

        if lines:
            replace_with_bullets(shape, lines, font_size_pt=font_size_pt)
            per["written"] = len(lines)
        else:
            per["notes"] = "no matching insights for provided categories"

        status[slide_key] = per

    return status

# =============================================================================
# Debug helpers (useful during template development)
# =============================================================================

def dump_placeholders(slide) -> List[Tuple[int, str, str, bool, str]]:
    """
    Inspect placeholders on a slide.
    Returns
    -------
    list of tuples
        (placeholder_idx, type_name, shape_name, text_capable, text_preview)
    """
    out = []
    for sh in slide.shapes:
        try:
            if getattr(sh, "is_placeholder", False):
                t = sh.placeholder_format.type
                tname = PP_PLACEHOLDER(t).name if isinstance(t, int) else str(t)
                text_capable = hasattr(sh, "text_frame")
                preview = ""
                if text_capable:
                    preview = (sh.text_frame.text or "").strip().replace("\n", "\\n")[:80]
                out.append((
                    getattr(sh.placeholder_format, "idx", -1),
                    tname, getattr(sh, "name", ""), text_capable, preview
                ))
        except Exception:
            # Keep going even if a shape misbehaves
            pass
    return out
