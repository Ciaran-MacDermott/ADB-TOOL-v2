"""
ppt_builder.py - FIXED
==============
Changes:
- Reverted bar spacing to original (gap_width=97, overlap=-27)
- Data labels remain OUTSIDE_END but no longer bold
- Charts now respect order of df_list (not alphabetically sorted)
"""
import os
import numpy as np
from itertools import accumulate
import pandas as pd
from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import (
    XL_CHART_TYPE, XL_TICK_LABEL_POSITION,
    XL_TICK_MARK, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
)
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor

from datetime import datetime
import re

from typing import Literal
from acc_deck_pkg.yoy_transformers import excel_round
from pptx import Presentation

# ============================================================
#  GLOBAL MAPPINGS (slide layouts, placeholders, etc.)
# ============================================================

SLIDE_LAYOUTS = {
    "intro": 0,
    "appendix": 8,
    "overlap": 6,
    "tot": 7,
    "all_other": 4,
    "misses": 5,
    "hits": 5,
    "category": 5  # Added for new category slides
}

# If a slide needs 1–4 charts (remainder), pick a layout that has exactly that many.
CHART_LAYOUT_DICT = {
    1: 2, 2: 3, 3: 5, 4: 4
}

# For a slide with N charts, these are the placeholder indices used as chart headers.
CHART_HEADER_PLACEHOLDERS = {
    1: [35],
    2: [35, 44],
    3: [35, 36, 37],
    4: [45, 56, 59, 58]
}

# For a slide with N charts, this placeholder holds the small italic footer text.
FOOTER_PLACEHOLDERS = {
    1: 37, 2: 37, 3: 34, 4: 37
}


# ============================================================
#  SLIDE CREATION HELPERS
# ============================================================


def add_slide(prs, slide_type):
    """Add a single slide of given `slide_type` (uses SLIDE_LAYOUTS mapping)."""
    layout_idx = SLIDE_LAYOUTS[slide_type]
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    return slide


def add_slides(prs, df_list, slide_type):
    """
    Add enough slides to hold all DataFrames for this `slide_type`.
    Packing rules (how many charts per slide):
      - "tot":       1 per slide (usually just one)
      - "all_other": 4 per slide (DEPRECATED - not used anymore)
      - "hits":      3 per slide (DEPRECATED - not used anymore)
      - "misses":    3 per slide (DEPRECATED - not used anymore)
      - "category":  3 per slide (NEW - main usage)
    """
    num_charts = len(df_list)

    def calc_slides(num_charts, slide_type):
        if slide_type == "tot":
            return num_charts, 0
        elif slide_type == "all_other":
            full, rem = divmod(num_charts, 4)
            return full, rem
        elif slide_type in ("hits", "misses", "category"):
            # Always 3 charts per slide
            full, rem = divmod(num_charts, 3)
            return full, rem
        else:
            return 0, 0

    full_slides, remainder = calc_slides(num_charts, slide_type)
    slide_list = []

    # Full slides
    for _ in range(full_slides):
        slide = add_slide(prs, slide_type)
        slide_list.append(slide)

    # Remainder slide (if needed)
    if remainder:
        layout_idx = CHART_LAYOUT_DICT[remainder]
        last_slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        slide_list.append(last_slide)

    return slide_list


# ============================================================
#  CHART CREATION
# ============================================================


def add_charts(prs, df_list, slide_list, metric_order=None, chart_footer=None,
               chart_subtitle="(YoY % Change)",
               header_font_size=21, header_underline=False):
    """
    Insert charts into the chart placeholders on each slide.
    NOW RESPECTS THE ORDER OF df_list (no sorting).

    header_font_size: font size in points for the per-chart header (default 21).
    header_underline: if True, the header is underlined (default False).
    """
    # REMOVED: df_list = sorted(df_list, key=lambda df: df.columns.name or "")
    # Charts now appear in the order they're passed in

    if chart_footer is None:
        chart_footer = "Smaller segment volume may shift categories between actuals updates"

    placeholder_groups = []
    header_groups = []

    # Discover chart and header placeholders on each slide
    for slide in slide_list:
        chart_phs = sorted(
            [s for s in slide.placeholders if s.name.startswith("Chart")],
            key=lambda s: s.left
        )
        placeholder_groups.append(chart_phs)

        header_idxs = CHART_HEADER_PLACEHOLDERS[len(chart_phs)]
        header_phs = sorted(
            [slide.placeholders[i] for i in header_idxs],
            key=lambda s: s.left
        )
        header_groups.append(header_phs)

    # Partition the df_list into chunks matching each slide's chart count
    df_splits = [len(p) for p in placeholder_groups]
    split_df_groups = [df_list[x - y: x] for x, y in zip(accumulate(df_splits), df_splits)]

    all_charts = []
    for i, slide in enumerate(slide_list):
        chart_placeholders = placeholder_groups[i]
        header_placeholders = header_groups[i]

        # Small-print footer on chart slides
        footer_idx = FOOTER_PLACEHOLDERS[len(chart_placeholders)]
        foot_frame = slide.placeholders[footer_idx].text_frame
        foot_frame.text = chart_footer
        for paragraph in foot_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(10)
                run.font.italic = True
                run.font.name = "Roboto Condensed"

        charts_for_slide = []
        for j, chart_ph in enumerate(chart_placeholders):
            df = split_df_groups[i][j]
            header_ph = header_placeholders[j]

            _add_chart_header(header_ph, df.columns.name, subtitle=chart_subtitle,
                              name_font_size=header_font_size,
                              name_underline=header_underline)
            chart = _create_clustered_chart(chart_ph, df, metric_order=metric_order)
            charts_for_slide.append(chart)

        all_charts.append(charts_for_slide)

    return all_charts


def _add_chart_header(header_ph, name, subtitle="(YoY % Change)",
                      name_font_size=21, name_underline=False):
    """Format and insert the two-line header above a chart. If subtitle is
    empty/None, the second paragraph is skipped — useful when the slide has
    a shared subtitle elsewhere (e.g. middle text box) and per-chart
    duplication is unwanted.

    name_font_size: font size in points for the chart header name (default 21).
    name_underline: if True, applies a single underline to the header name
                    (default False)."""
    header_frame = header_ph.text_frame
    header_frame.word_wrap = len(name) >= 25
    header_frame.vertical_anchor = MSO_ANCHOR.TOP
    header_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    p1 = header_frame.paragraphs[0]
    p1.space_after = Pt(0.5)
    r1 = p1.add_run()
    r1.text = name
    r1.font.size = Pt(name_font_size)
    r1.font.bold = True
    r1.font.underline = bool(name_underline)
    r1.font.name = "Roboto Condensed"
    r1.font.color.rgb = RGBColor(0, 0, 0)

    if subtitle:
        p2 = header_frame.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(15)
        r2.font.bold = False
        r2.font.name = "Roboto Condensed"
        r2.font.color.rgb = RGBColor(0, 0, 0)


# ================== CLUSTERED (GROUPED) BAR LOGIC ==================


def _create_clustered_chart(chart_ph, df, metric_order=None):
    """
    Build a clustered column chart (Forecast vs Actual) across metrics.
    """
    pv = df.pivot(index="metric", columns="f_or_a", values="YoY").copy()

    # Ensure both series exist
    for col in ("Forecast", "Actual"):
        if col not in pv.columns:
            pv[col] = 0.0
    pv = pv[["Forecast", "Actual"]].fillna(0.0)

    # x-axis order: use provided order, fall back to default, then to DataFrame order
    if metric_order is None:
        metric_order = ["Units", "ASP", "Dollars"]
    order = [m for m in metric_order if m in pv.index]
    if not order:
        order = list(pv.index)
    pv = pv.loc[order]

    chart_data = CategoryChartData()
    chart_data.categories = list(pv.index)
    chart_data.add_series("Forecast", list(pv["Forecast"]))
    chart_data.add_series("Actual", list(pv["Actual"]))

    graphic_frame = chart_ph.insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data)
    chart = graphic_frame.chart

    # Color identity
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = RGBColor(239, 42, 121)  # Forecast: pink
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = RGBColor(78, 16, 111)  # Actual: purple
    return chart


# ============================================================
#  TITLE AND FOOTER TEXT FORMATTING
# ============================================================

def _format_month_year(val, fallback_dt: datetime | None = None) -> str:
    """
    Return a 'Month Year' label from various inputs.
    """
    # datetime-like
    if hasattr(val, "strftime"):
        return val.strftime("%B %Y")

    # string cases
    if isinstance(val, str):
        s = val.strip()

        # 04/2025 or 4/2025
        m = re.fullmatch(r"(\d{1,2})/(\d{4})", s)
        if m:
            month = int(m.group(1))
            year = int(m.group(2))
            try:
                return datetime(year, month, 1).strftime("%B %Y")
            except ValueError:
                pass

        # 2025-04 or 2025-4
        m = re.fullmatch(r"(\d{4})-(\d{1,2})", s)
        if m:
            year = int(m.group(1))
            month = int(m.group(2))
            try:
                return datetime(year, month, 1).strftime("%B %Y")
            except ValueError:
                pass

        # Already Month Year
        for fmt in ("%B %Y", "%b %Y"):
            try:
                dt = datetime.strptime(s, fmt)
                return dt.strftime("%B %Y")
            except ValueError:
                continue

        return s

    # fallback
    dt = fallback_dt or datetime.now()
    return dt.strftime("%B %Y")


def apply_titles(prs, config, type_dict_pres, today, default_type: str | None = None):
    """
    Apply titles and subtext per slide type using CONFIG.
    """
    from pptx.util import Pt

    release_raw = config.get("release")
    release_label = _format_month_year(release_raw, fallback_dt=today)

    input_level1 = config["input_level1"]
    input_year = config["input_year"]
    input_quarter = config["input_quarter"]

    title_dict = {
        "intro": f"Q{input_quarter} {input_year} Forecast vs. Actuals",
        "tot": f"{input_level1}: Total Industry Q{input_quarter} {input_year}",
        "hits": f"Q{input_quarter} {input_year} Best Performing Categories",
        "misses": f"Q{input_quarter} {input_year} Biggest Miss Categories",
        "appendix": "Appendix",
        "overlap": f"Q{input_quarter} {input_year} Forecast Note:",
        "all_other": f"All Other {input_level1} Categories",
        "category": f"{input_level1} Categories - Q{input_quarter} {input_year}",
    }

    for idx, slide in enumerate(prs.slides):
        type_key = type_dict_pres.get(idx, default_type)
        if not type_key:
            continue
        try:
            title_frame = slide.placeholders[0].text_frame
        except Exception:
            continue

        prev_same = idx > 0 and type_dict_pres.get(idx - 1) == type_key
        title_frame.text = f"{title_dict[type_key]}{'' if prev_same else ''}"

        # Footer continuation hint
        try:
            next_same = type_dict_pres.get(idx + 1) == type_key
            if next_same:
                foot_frame = slide.placeholders[38].text_frame
                foot_frame.text = "Continued on Next Page"
                for run in foot_frame.paragraphs[0].runs:
                    run.font.size = Pt(10)
                    run.font.italic = True
        except Exception:
            pass

    # Intro slide extras
    if type_dict_pres.get(0) == "intro":
        try:
            intro = prs.slides[0]
            intro.placeholders[12].text_frame.text = f"Future of™ insights for {input_level1}"
            intro.placeholders[11].text_frame.text = f"Report Prepared: {today.strftime('%B %Y')}"
        except Exception:
            pass


# ============================================================
#  FINAL CHART FORMATTING (REVERTED TO ORIGINAL STYLE)
# ============================================================


def _normalize_pres_charts(pres_charts):
    """
    Normalize different shapes of 'pres_charts' into list[list[chart]].
    """
    if pres_charts and hasattr(pres_charts[0], "value_axis"):
        return [pres_charts]

    norm = []
    for layer in pres_charts:
        if not layer:
            continue
        if hasattr(layer, "value_axis"):
            norm.append([layer])
        elif isinstance(layer, list):
            charts = []
            for item in layer:
                if hasattr(item, "value_axis"):
                    charts.append(item)
                elif isinstance(item, list):
                    charts.extend([x for x in item if hasattr(x, "value_axis")])
            if charts:
                norm.append(charts)
    return norm


def format_charts(prs, pres_charts, pres_dfs, min_axis_range=15.0, use_fixed_scale=False, fixed_min=-30.0,
                  fixed_max=30.0):
    """
    Apply consistent visual formatting across all charts in the deck.
    REVERTED: Bar spacing back to original for cleaner look.
    KEPT: OUTSIDE_END positioning, minimum axis range.
    """
    pres_charts = _normalize_pres_charts(pres_charts)
    charts = [c for slide_charts in pres_charts for c in slide_charts]

    # Calculate axis range
    if use_fixed_scale:
        min_val, max_val = fixed_min, fixed_max
        print(f"Using fixed axis scale: {min_val:.1f} to {max_val:.1f}")
    else:
        try:
            max_val = float(np.nanmax([df["YoY"].max() for df in pres_dfs]))
            min_val = float(np.nanmin([df["YoY"].min() for df in pres_dfs]))
            if not np.isfinite(max_val) or not np.isfinite(min_val):
                raise ValueError

            # Ensure minimum range to make small changes visible
            data_range = max_val - min_val
            if data_range < min_axis_range:
                # Expand range symmetrically around the center
                center = (max_val + min_val) / 2
                max_val = center + min_axis_range / 2
                min_val = center - min_axis_range / 2
                print(f"Expanded axis range to minimum {min_axis_range:.1f} points: {min_val:.1f} to {max_val:.1f}")
            else:
                print(f"Using data-driven axis scale: {min_val:.1f} to {max_val:.1f}")

        except Exception as e:
            max_val, min_val = 100.0, -100.0
            print(f"Warning: Could not calculate axis range, using defaults: {e}")

    for chart in charts:
        # Set axis range
        chart.value_axis.maximum_scale = max_val
        chart.value_axis.minimum_scale = min_val

        # Remove tick marks
        chart.value_axis.minor_tick_mark = XL_TICK_MARK.NONE
        chart.category_axis.minor_tick_mark = XL_TICK_MARK.NONE
        chart.value_axis.major_tick_mark = XL_TICK_MARK.NONE
        chart.category_axis.major_tick_mark = XL_TICK_MARK.NONE
        chart.value_axis.visible = False

        # Category axis positioning
        chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
        chart.category_axis.tick_labels.offset = 500

        # Category axis label formatting
        chart.category_axis.tick_labels.font.size = Pt(12)
        chart.category_axis.tick_labels.font.name = "Roboto Condensed"

        # Force horizontal rotation (no diagonal) via XML
        _ns = {
            'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        }
        axis_el = chart.category_axis._element
        txPr = axis_el.find('c:txPr', namespaces=_ns)
        if txPr is None:
            txPr = etree.SubElement(
                axis_el,
                '{http://schemas.openxmlformats.org/drawingml/2006/chart}txPr',
            )
        bodyPr = txPr.find('a:bodyPr', namespaces=_ns)
        if bodyPr is None:
            bodyPr = etree.SubElement(
                txPr,
                '{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr',
            )
        bodyPr.set('rot', '0')
        bodyPr.set('vert', 'horz')

        # Remove gridlines
        chart.value_axis.has_major_gridlines = False
        chart.category_axis.has_major_gridlines = False
        chart.value_axis.has_minor_gridlines = False
        chart.category_axis.has_minor_gridlines = False

        # Legend
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.font.name = "Roboto Condensed"
        chart.legend.font.size = Pt(14)
        chart.legend.include_in_layout = False

        # REVERTED: Original bar spacing for cleaner look
        for plot in chart.plots:
            plot.gap_width = 97  # BACK TO ORIGINAL (was 50)
            plot.overlap = -27  # BACK TO ORIGINAL (was -10)

            for series in plot.series:
                series.invert_if_negative = False

                # Enable data labels at series level first
                series.has_data_labels = True
                lbl = series.data_labels
                lbl.show_value = True
                lbl.show_category_name = False
                lbl.show_series_name = False
                lbl.font.size = Pt(12)
                lbl.font.name = "Roboto Condensed"
                lbl.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

                # Override each point's label text with formatted value
                for point, val in zip(series.points, series.values):
                    dl = point.data_label
                    dl.has_text_frame = True
                    dl.text_frame.text = f"{excel_round(val, decimals=1)}%"
                    for run in dl.text_frame.paragraphs[0].runs:
                        run.font.size = Pt(12)
                        run.font.name = "Roboto Condensed"


def set_text_in_placeholder_31(
        prs: Presentation,
        text: str,
        mode: Literal["replace", "append"] = "replace",
) -> int:
    """
    Set or append text for placeholder index 31 across all slides.
    """
    updated = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                pf = shape.placeholder_format
            except (AttributeError, ValueError):
                continue

            if pf and pf.idx == 31:
                if not hasattr(shape, "text_frame"):
                    continue

                tf = shape.text_frame
                if mode == "replace":
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = text
                    for run in p.runs:
                        run.font.name = "Roboto Condensed"
                else:  # append
                    p = tf.add_paragraph()
                    p.text = text
                    for run in p.runs:
                        run.font.name = "Roboto Condensed"
                updated += 1

    return updated


def remove_empty_placeholders(prs: Presentation, skip_indices: set | None = None) -> int:
    """
    Remove empty placeholders from all slides to make deck more client-ready.

    Args:
        prs: Presentation object.
        skip_indices: Set of placeholder indices to preserve even if empty
                      (e.g. {14} to keep subheader for LLM population).

    Returns the count of placeholders removed.
    """
    if skip_indices is None:
        skip_indices = set()

    removed = 0

    for slide in prs.slides:
        shapes_to_remove = []

        for shape in slide.shapes:
            # Check if it's a placeholder
            try:
                pf = shape.placeholder_format
            except (AttributeError, ValueError):
                continue
            if not pf:
                continue

            # Skip preserved indices
            if pf.idx in skip_indices:
                continue

            # Check if it has a text frame
            if not hasattr(shape, "text_frame"):
                continue

            # Check if empty (no text content)
            tf = shape.text_frame
            all_text = "".join(p.text for p in tf.paragraphs).strip()

            if not all_text:
                shapes_to_remove.append(shape)

        # Remove empty placeholders
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
            removed += 1

    return removed