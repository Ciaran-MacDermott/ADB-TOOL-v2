#!/usr/bin/env python3
"""
pipeline.py
NPD Foodservice Pipeline
Orchestrates data loading, preparation, and PowerPoint generation.
"""

import os
import sys
import json
import shutil
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.dml.color import RGBColor

# Add project root so acc_deck_pkg is importable
PROJECT_ROOT = Path(__file__).resolve().parent.parent
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from acc_deck_fs_pkg.data_analysis import load_and_merge_data
from acc_deck_fs_pkg.data_prep import (
    prepare_data, build_total_chart_df, build_segment_chart_df,
    build_prefix_split_chart_dfs, build_food_bev_chart_df, METRIC_ORDER,
    FOOD_DRINK_MAP,
)
from acc_deck_pkg import ppt_builder
from acc_deck_fs_pkg.llm_insights import generate_all_insights, apply_insights_to_slides

# ============================================================================
# PATHS
# ============================================================================

# _BASE_DIR: package directory containing bundled assets (Templates/, Images/).
# _WRITABLE_DIR: user-writable data at runtime — CSVs land here.
#   Override via env FS_WRITABLE_DIR (Streamlit sets this to a per-run tempdir).
_BASE_DIR     = Path(getattr(sys, '_MEIPASS', Path(__file__).resolve().parent))
_WRITABLE_DIR = Path(os.environ.get("FS_WRITABLE_DIR", str(_BASE_DIR)))

DATA_DIR = Path(__file__).parent   # acc_deck_fs_pkg/ — for config.json and module-relative reads

TEMPLATE_PATH = _BASE_DIR     / "Templates" / "template.pptx"
INTRO_IMAGE   = _BASE_DIR     / "Images"    / "food service title img.jpg"
FORECAST_CSV  = _WRITABLE_DIR / "Dashboard API data" / "forecast_full.csv"
ACTUALS_CSV   = _WRITABLE_DIR / "Dashboard API data" / "actuals_full.csv"

# ============================================================================
# PIPELINE CONFIGURATION
# ============================================================================

# Load config from shared config.json (model params, pipeline thresholds — never API keys)
_CONFIG_PATH = DATA_DIR / "config.json"
_config      = json.loads(_CONFIG_PATH.read_text(encoding="utf-8")) if _CONFIG_PATH.exists() else {}
_pipeline_section = _config.get("pipeline", {})

PIPELINE_CONFIG = {
    'input_year':        _pipeline_section.get("input_year", 2025),
    'input_quarter':     _pipeline_section.get("input_quarter", 4),
    'comp_year':         _pipeline_section.get("comp_year", None),     # None → input_year - 1
    'comp_quarter':      _pipeline_section.get("comp_quarter", None),  # None → same as input_quarter
    'metrics':           _pipeline_section.get("metrics", ['dollars', 'units', 'asp']),
    'groq_api_key':      os.environ.get("GROQ_API_KEY", ""),
    'moonshot_api_key':  os.environ.get("MOONSHOT_API_KEY", ""),
    'generate_insights': _pipeline_section.get("generate_insights", True),
    'output_path':       None,   # Set by caller; None → auto-generate from DATA_DIR
    'industry_id_key':   'food-service',
    'project_display':   None,   # Set by caller (e.g. "Food Service Canada")
}

# Slide 2: segment categories per industry.
# Short display names are fine — _find_level2 maps them via LABEL_ALIASES.
# Add new industry keys here to customise without touching main().
SEGMENT_CATEGORIES = {
    'food-service': [
        'Full Service Restaurants',
        'Quick Service Restaurants',
        # Display label for the API's "Fast Casual" segment
        # (resolved via LABEL_ALIASES in data_prep.py).
        'Fast Casual Restaurants',
    ],
    'food-service-canada': [
        # Display labels — "Foodservice" suffix dropped per client convention.
        # _find_level2() resolves these to the actual API level2 values via
        # LABEL_ALIASES ("Total Commercial" → "Total Commercial Foodservice",
        # "Retail" → "Retail Foodservice").
        'Total Commercial',
        'Quick Service Restaurants',
        'Full Service Restaurants',
        'Retail',
    ],
    'food-service-australia': [
        'Full Service Restaurants',
        'Quick Service Restaurants',
        'Retail',
    ],
    # Default fallback (used for any industry not listed above)
    '_default': [
        'Full Service Restaurants',
        'Quick Service Restaurants',
    ],
}

# Slides 3 + 4 + 5: chart layout per industry.
# split_dayparts / split_service_modes = True  → two charts (QSR / FSR side-by-side, "category" slide)
#                                       = False → one combined chart ("tot" slide)
# top_food / top_drink control the food & beverage slide item counts.
# Add a new key here to customise a market without touching main().
CHART_CONFIG = {
    '_default': {
        # US foodservice.
        'total_level2': 'Total Restaurants',
        'cut_term':     'Restaurant Performance',   # used in slide 3/4/5 titles
        'split_dayparts': True,
        'split_service_modes': True,
        # US uses "Dinner" (matches the API's QSR Dinner / FSR Dinner directly,
        # and matches the previous-wave US template prose).
        'daypart_items': ['Morning Meal', 'Lunch', 'Dinner', 'P.M. Snack'],
        'service_mode_items': {
            'QSR': ['On-Premises', 'Carry-Out', 'Drive-Thru', 'Delivery'],
            'FSR': ['On-Premises', 'Carry-Out', 'Delivery'],
        },
        'top_food': 5,
        'top_drink': 3,
    },
    'food-service-canada': {
        'total_level2':         'Total Restaurants',
        # CA slide 3 uses an umbrella-prefix override: "Total Commercial Foodservice Q<n> <yyyy>"
        'segments_chart_title': 'Total Commercial Foodservice',
        # Slides 4/5 still use "Restaurant Performance" via cut_term.
        'cut_term':             'Restaurant Performance',
        'split_dayparts': True,
        'split_service_modes': True,
        'daypart_items': ['Morning Meal', 'Lunch', 'Supper', 'P.M. Snack'],
        'service_mode_items': {
            'QSR': ['On-Premises', 'Carry-Out', 'Drive-Thru', 'Delivery'],
            'FSR': ['On-Premises', 'Carry-Out', 'Delivery'],
        },
        'top_food': 5,
        'top_drink': 3,
    },
    'food-service-australia': {
        # AUS slides 3/4/5 share "Commercial Foodservice by ..." cut_term phrasing.
        'total_level2': 'Total Commercial Foodservice',
        'cut_term':     'Commercial Foodservice',
        'split_dayparts': False,
        'split_service_modes': False,
        'daypart_items': ['Morning Meal', 'Lunch', 'Dinner', 'P.M. Snack'],
        'service_mode_items': ['On-Premises', 'Carry Out', 'Drive-Thru', 'Delivery/Pickup'],
        'top_food': 3,
        'top_drink': 3,
    },
}


def _derive_project_label(df):
    """Return the display project name, preferring the API label from the GUI."""
    # GUI sets project_display from the industries API label (e.g. "Food Service Canada")
    from_config = PIPELINE_CONFIG.get('project_display')
    if from_config:
        return from_config
    # Fallback: derive from the project column in the data
    raw = df['project'].iloc[0] if 'project' in df.columns else 'foodservice'
    return raw.replace('-', ' ').title()


def _geo_suffix() -> str:
    """
    Return ' U.S.' only for the base US food-service industry.
    Regional variants (food-service-canada, food-service-uk, etc.) already
    include the country in their label, so no suffix is needed.
    """
    return ' U.S.' if PIPELINE_CONFIG.get('industry_id_key', '') == 'food-service' else ''


def _set_slide_title(slide, text, font_size=36):
    """Set the title placeholder (idx 0) on a slide."""
    try:
        tf = slide.placeholders[0].text_frame
        tf.text = text
        for p in tf.paragraphs:
            for run in p.runs:
                run.font.size = Pt(font_size)
                run.font.name = "Poppins"
    except (KeyError, AttributeError):
        pass


def _add_food_bev_separator(slide, num_food, num_total):
    """Add a faint grey dashed vertical line between food and beverage categories."""
    # Find the chart shape on this slide
    chart_shape = None
    for shape in slide.shapes:
        if shape.has_chart:
            chart_shape = shape
            break
    if chart_shape is None:
        return

    # Estimate plot area within the chart (auto-layout defaults)
    chart_left = chart_shape.left
    chart_top = chart_shape.top
    chart_w = chart_shape.width
    chart_h = chart_shape.height

    plot_left_pct = 0.07   # value-axis label margin
    plot_right_pct = 0.98
    plot_left = chart_left + int(chart_w * plot_left_pct)
    plot_width = int(chart_w * (plot_right_pct - plot_left_pct))

    sep_x = plot_left + int(plot_width * num_food / num_total)
    line_top = chart_top + int(chart_h * 0.12)   # below legend
    line_bot = chart_top + int(chart_h * 0.90)   # above category labels

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        sep_x, line_top,
        sep_x, line_bot,
    )
    connector.line.color.rgb = RGBColor(180, 180, 180)
    connector.line.width = Pt(0.75)
    connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH


def _set_shared_chart_title(slide, line1, line2, placeholder_idx=40):
    """Populate the shared chart title placeholder (idx 37) on 2-chart slides.
    If line2 is empty/None, the subtitle paragraph is skipped so no blank
    line takes up vertical space below the title."""
    try:
        ph = slide.placeholders[placeholder_idx]
        tf = ph.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = line1
        r1.font.size = Pt(21)
        r1.font.bold = True
        r1.font.name = "Roboto Condensed"
        r1.font.color.rgb = RGBColor(0, 0, 0)

        if line2:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = line2
            r2.font.size = Pt(15)
            r2.font.bold = False
            r2.font.name = "Roboto Condensed"
            r2.font.color.rgb = RGBColor(0, 0, 0)
    except (KeyError, AttributeError):
        pass


def main(df_forecast=None, df_actuals=None):
    """
    Build the forecast accuracy deck.

    Args:
        df_forecast: Forecast DataFrame from extract_data() — when supplied the
                     CSV fallback is skipped entirely (preferred live path).
        df_actuals:  Actuals DataFrame from extract_data() — same as above.
    """
    cfg = PIPELINE_CONFIG
    input_year = cfg['input_year']
    input_quarter = cfg['input_quarter']
    metrics = cfg['metrics']

    print("\n" + "=" * 60)
    print("NPD FOODSERVICE PIPELINE")
    print("=" * 60)

    # === Stage 1: Load & merge data ===
    # When DataFrames are passed in (live extract path) they are used directly.
    # Falls back to reading the saved CSVs only when called standalone.
    print("\n[1/5] Loading data...")
    df = load_and_merge_data(
        forecast_path=str(FORECAST_CSV),
        actuals_path=str(ACTUALS_CSV),
        metrics=metrics,
        df_forecast=df_forecast,
        df_actuals=df_actuals,
    )
    print(f"      Merged: {len(df):,} rows, columns: {list(df.columns)}")

    project_label = _derive_project_label(df)

    # === Stage 2: Prepare chart DataFrames ===
    print("\n[2/5] Preparing chart data...")

    # Common chart header/subtitle strings
    chart_title = f"Q{input_quarter} Forecasts vs Actuals"
    traffic_subtitle = "(Traffic YoY % Change)"

    # Resolve per-industry chart config early (needed for Slide 1 total_level2)
    industry_key = cfg.get('industry_id_key', 'food-service')
    chart_cfg = CHART_CONFIG.get(industry_key, CHART_CONFIG['_default'])
    total_level2 = chart_cfg.get('total_level2') or 'Total Restaurants'

    # Slide 1: Total — single chart, Forecast vs Actual YoY%
    total_chart_df = build_total_chart_df(
        df, input_year, input_quarter,
        chart_label=chart_title,
        total_level2=total_level2,
    )
    print(f"      Slide 1 (Total): {len(total_chart_df)} rows")

    # Slide 2: Restaurant segments — industry-specific category list
    segment_categories = SEGMENT_CATEGORIES.get(industry_key, SEGMENT_CATEGORIES['_default'])
    print(f"      Slide 2 categories ({industry_key}): {segment_categories}")
    segment_chart_df = build_segment_chart_df(
        df, input_year, input_quarter,
        categories=segment_categories,
        metric='units',
        chart_label=chart_title,
    )
    print(f"      Slide 2 (Segments): {len(segment_chart_df)} rows")

    # Slide 3: Dayparts
    daypart_items = chart_cfg['daypart_items']
    if chart_cfg['split_dayparts']:
        # Two charts: QSR left, FSR right
        daypart_dfs = build_prefix_split_chart_dfs(
            df, input_year, input_quarter,
            prefix_map={
                'QSR': 'Quick Service Restaurants',
                'FSR': 'Full-Service Restaurants',
            },
            items=daypart_items,
            metric='units',
        )
        daypart_slide_type = "category"
    else:
        # Single combined chart (no QSR/FSR prefix in data)
        combined_dp = build_segment_chart_df(
            df, input_year, input_quarter,
            categories=daypart_items,
            metric='units',
            chart_label=chart_title,
        )
        daypart_dfs = [combined_dp]
        daypart_slide_type = "tot"
    print(f"      Slide 3 (Dayparts): {len(daypart_dfs)} chart(s), "
          f"split={chart_cfg['split_dayparts']}")

    # Slide 4: Service Modes
    svc_items = chart_cfg['service_mode_items']
    if chart_cfg['split_service_modes']:
        # Two charts: QSR left, FSR right
        service_mode_dfs = build_prefix_split_chart_dfs(
            df, input_year, input_quarter,
            prefix_map={
                'QSR': 'Quick Service Restaurants',
                'FSR': 'Full-Service Restaurants',
            },
            items=svc_items,
            metric='units',
        )
        svc_slide_type = "category"
    else:
        # Single combined chart
        combined_svc = build_segment_chart_df(
            df, input_year, input_quarter,
            categories=svc_items,
            metric='units',
            chart_label=chart_title,
        )
        service_mode_dfs = [combined_svc]
        svc_slide_type = "tot"
    print(f"      Slide 4 (Service Modes): {len(service_mode_dfs)} chart(s), "
          f"split={chart_cfg['split_service_modes']}")

    # Slide 5: Food & Beverage
    servings_subtitle = "(Servings YoY % Change)"
    food_bev_df, food_bev_order = build_food_bev_chart_df(
        df, input_year, input_quarter,
        top_food=chart_cfg.get('top_food', 5),
        top_drink=chart_cfg.get('top_drink', 3),
        metric='units',
        chart_label=chart_title,
    )
    print(f"      Slide 5 (Food & Bev): {len(food_bev_order)} categories "
          f"({food_bev_order})")

    # === Stage 3: Setup presentation ===
    print("\n[3/5] Setting up presentation...")
    prs = Presentation(str(TEMPLATE_PATH))
    type_map = {}
    today = datetime.now()

    # --- Intro slide ---
    intro_slide = ppt_builder.add_slide(prs, "intro")
    type_map[prs.slides.index(intro_slide)] = "intro"
    _set_slide_title(intro_slide, f"Q{input_quarter} {input_year} Forecast vs. Actuals", font_size=56)
    try:
        intro_slide.placeholders[12].text_frame.text = (
            f"Future of\u2122 insights for {project_label}{_geo_suffix()}"
        )
        intro_slide.placeholders[11].text_frame.text = (
            f"Report Prepared: {today.strftime('%B %Y')}"
        )
    except (KeyError, AttributeError):
        pass

    # Add title image — fills the picture placeholder (idx 14)
    try:
        intro_slide.placeholders[14].insert_picture(str(INTRO_IMAGE))
    except (KeyError, AttributeError):
        pass

    # --- Slide 1: Total Restaurants ---
    total_slides = ppt_builder.add_slides(prs, [total_chart_df], "tot")
    for s in total_slides:
        type_map[prs.slides.index(s)] = "tot"

    total_charts = ppt_builder.add_charts(
        prs,
        [total_chart_df],
        total_slides,
        metric_order=METRIC_ORDER,
        chart_footer="",
    )

    # Slide 2 title: derived from total_level2 — "Total Restaurants Q4 2025" for
    # US/CA, "Total Commercial Foodservice Q4 2025" for AUS.
    _set_slide_title(
        total_slides[0],
        f"{total_level2} Q{input_quarter} {input_year}",
    )

    # --- Slide 2: Restaurant Segments (Traffic) ---
    segment_slides = ppt_builder.add_slides(prs, [segment_chart_df], "tot")
    for s in segment_slides:
        type_map[prs.slides.index(s)] = "tot"

    segment_charts = ppt_builder.add_charts(
        prs,
        [segment_chart_df],
        segment_slides,
        metric_order=segment_categories,
        chart_footer="",
        chart_subtitle=traffic_subtitle,
    )

    # Slide 3 title:
    #   - When CHART_CONFIG defines a 'segments_chart_title' override the slide
    #     title becomes "<override> Q<n> <yyyy>" (CA uses this, e.g.
    #     "Total Commercial Foodservice Q4 2025").
    #   - Otherwise falls back to "Q<n> <yyyy> <cut_term> by Segments"
    #     (US: "Q4 2025 Restaurant Performance by Segments";
    #      AUS: "Q4 2025 Commercial Foodservice by Segments").
    cut_term = chart_cfg.get('cut_term', 'Restaurant Performance')
    segments_title_override = chart_cfg.get('segments_chart_title')
    if segments_title_override:
        slide3_title = f"{segments_title_override} Q{input_quarter} {input_year}"
    else:
        slide3_title = f"Q{input_quarter} {input_year} {cut_term} by Segments"
    _set_slide_title(segment_slides[0], slide3_title)

    # --- Slide 3: Dayparts ---
    daypart_slides = ppt_builder.add_slides(prs, daypart_dfs, daypart_slide_type)
    for s in daypart_slides:
        type_map[prs.slides.index(s)] = daypart_slide_type

    # When dayparts are split (US/CA) the per-chart subtitle is dropped because
    # the middle text box already shows "(Traffic YoY % Change)". When NOT split
    # (AUS — single combined chart, no middle text box) the subtitle goes on the
    # chart itself so the audience still sees the metric label.
    daypart_chart_subtitle = "" if chart_cfg['split_dayparts'] else traffic_subtitle
    daypart_charts = ppt_builder.add_charts(
        prs,
        daypart_dfs,
        daypart_slides,
        metric_order=daypart_items,
        chart_footer="",
        chart_subtitle=daypart_chart_subtitle,
        header_font_size=16,      # smaller than the default 21 per Slide 4 (Dayparts) spec
        header_underline=True,    # underline the QSR / FSR chart header
    )

    _set_slide_title(
        daypart_slides[0],
        f"Q{input_quarter} {input_year} {cut_term} by Dayparts",
    )
    if chart_cfg['split_dayparts']:
        # Middle text box on Slide 4 (Dayparts) keeps both the title AND the
        # "(Traffic YoY % Change)" subtitle — client convention. The duplicate
        # "(YoY % Change)" on each individual chart header is dropped via
        # chart_subtitle="" in the add_charts() call above.
        _set_shared_chart_title(daypart_slides[0], chart_title, traffic_subtitle)

    # --- Slide 4: Service Modes ---
    svc_slides = ppt_builder.add_slides(prs, service_mode_dfs, svc_slide_type)
    for s in svc_slides:
        type_map[prs.slides.index(s)] = svc_slide_type

    # Same split-vs-combined rationale as dayparts: drop chart subtitle when the
    # middle text box carries it (US/CA), keep it on the chart when AUS combines
    # everything into a single chart with no middle text box.
    svc_chart_subtitle = "" if chart_cfg['split_service_modes'] else traffic_subtitle
    svc_charts = ppt_builder.add_charts(
        prs,
        service_mode_dfs,
        svc_slides,
        chart_footer="",
        chart_subtitle=svc_chart_subtitle,
        header_font_size=16,      # smaller than the default 21 per Slide 5 (Service Modes) spec
        header_underline=True,    # underline the QSR / FSR chart header
    )

    _set_slide_title(
        svc_slides[0],
        f"Q{input_quarter} {input_year} {cut_term} by Service Mode",
        font_size=34,
    )
    if chart_cfg['split_service_modes']:
        # Middle text box on Slide 5 (Service Modes) keeps both the title AND
        # the "(Traffic YoY % Change)" subtitle — client convention. The
        # duplicate "(YoY % Change)" on each individual chart header is
        # dropped via chart_subtitle="" in the add_charts() call above.
        _set_shared_chart_title(svc_slides[0], chart_title, traffic_subtitle)

    # --- Slide 5: Food & Beverage ---
    fb_slides = ppt_builder.add_slides(prs, [food_bev_df], "tot")
    for s in fb_slides:
        type_map[prs.slides.index(s)] = "tot"

    # Slide 6 footnote: phrasing depends on whether top_food == top_drink.
    #   - US/CA (5 food / 3 drink) → "Top 5 food items and top 3 beverage items..."
    #   - AUS  (3 food / 3 drink)  → "Top 3 food/beverage are listed..."
    top_food = chart_cfg.get('top_food', 5)
    top_drink = chart_cfg.get('top_drink', 3)
    if top_food == top_drink:
        fb_footer = (
            f"Top {top_food} food/beverage are listed based on "
            f"Q{input_quarter} actual servings"
        )
    else:
        fb_footer = (
            f"Top {top_food} food items and top {top_drink} beverage items are listed "
            f"based on Q{input_quarter} actual servings"
        )

    fb_charts = ppt_builder.add_charts(
        prs,
        [food_bev_df],
        fb_slides,
        metric_order=food_bev_order,
        chart_footer=fb_footer,
        chart_subtitle=servings_subtitle,
    )

    _set_slide_title(
        fb_slides[0],
        f"Q{input_quarter} {input_year} Food & Beverage Performance",
    )

    # Add faint vertical separator between food and beverage categories
    _add_food_bev_separator(fb_slides[0], num_food=chart_cfg.get('top_food', 5), num_total=len(food_bev_order))

    # Format all charts together (consistent axis range across slides)
    all_charts = total_charts + segment_charts + daypart_charts + svc_charts + fb_charts
    all_dfs = [total_chart_df, segment_chart_df] + daypart_dfs + service_mode_dfs + [food_bev_df]
    ppt_builder.format_charts(prs, all_charts, all_dfs)

    # === Stage 4: Generate LLM insights ===
    if cfg.get('generate_insights', False) and (cfg.get('groq_api_key') or cfg.get('moonshot_api_key')):
        print("\n[4/5] Generating LLM insights...")
        insights = generate_all_insights(
            api_key=cfg['groq_api_key'],
            total_chart_df=total_chart_df,
            segment_chart_df=segment_chart_df,
            daypart_dfs=daypart_dfs,
            service_mode_dfs=service_mode_dfs,
            food_bev_df=food_bev_df,
            food_bev_map=FOOD_DRINK_MAP,
            industry_id_key=industry_key,
        )

        # Map slide names to slide objects for insight placement
        slide_map = {
            "total": total_slides[0],
            "segments": segment_slides[0],
            "dayparts": daypart_slides[0],
            "service_modes": svc_slides[0],
            "food_bev": fb_slides[0],
        }
        applied = apply_insights_to_slides(prs, insights, slide_map)
        print(f"      Applied {applied} insights to slides")
    else:
        print("\n[4/5] Skipping LLM insights (disabled or no API key)")

    # === Finalize: footer + cleanup ===
    print("\n[5/5] Finalizing...")
    ppt_builder.set_text_in_placeholder_31(
        prs,
        f"Future of\u2122 insights for {project_label}{_geo_suffix()} Forecast Accuracy - Q{input_quarter} {input_year}",
        mode="replace",
    )
    removed = ppt_builder.remove_empty_placeholders(prs, skip_indices={14})
    print(f"      Removed {removed} empty placeholders")

    # === Save ===
    # Use GUI-supplied path if set, otherwise fall back to auto-generated name
    gui_path = cfg.get('output_path')
    if gui_path:
        output_path = Path(gui_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
    else:
        output_path = DATA_DIR / f"Future_of_{project_label}_Accuracy_Deck_{input_year}_Q{input_quarter}.pptx"

    try:
        prs.save(str(output_path))
    except PermissionError:
        output_path = output_path.with_stem(output_path.stem + "_v2")
        prs.save(str(output_path))
    print(f"      Saved: {output_path}")

    # Save raw + merged data CSVs alongside the PowerPoint for analyst validation
    data_out_dir = output_path.parent
    stem = f"{project_label.replace(' ', '_')}_{input_year}_Q{input_quarter}"

    fc_src = FORECAST_CSV.resolve()
    ac_src = ACTUALS_CSV.resolve()

    if fc_src.exists():
        dest = data_out_dir / f"forecast_raw_{stem}.csv"
        shutil.copy2(str(fc_src), str(dest))
        print(f"      Saved: {dest.name}")
    if ac_src.exists():
        dest = data_out_dir / f"actuals_raw_{stem}.csv"
        shutil.copy2(str(ac_src), str(dest))
        print(f"      Saved: {dest.name}")

    merged_dest = data_out_dir / f"merged_data_{stem}.csv"
    df.to_csv(str(merged_dest), index=False)
    print(f"      Saved: {merged_dest.name}")

    print("\n" + "=" * 60)
    print("PIPELINE COMPLETE")
    print("=" * 60)

    return {"prs": prs, "df": df, "type_map": type_map}


def run_full_pipeline(prod_session=None, qa_session=None,
                      industry_id: str = 'food-service', extract: bool = True):
    """
    Connected pipeline: extract fresh data → save CSVs → build deck.

    Login is handled by the GUI before this is called. Sessions are passed
    in directly so no Selenium login occurs here.

    Args:
        prod_session: Authenticated requests.Session for Production (forecast).
        qa_session:   Authenticated requests.Session for QA (actuals).
        industry_id:  Industry ID (e.g. 'food-service', 'food-service-uk').
        extract:      If True, fetch fresh data before building the deck.
                      If False, assumes CSVs already exist.

    Returns:
        Result dict from main() (prs, df, type_map).
    """
    # Propagate industry choice into module-level config so main() picks the right
    # SEGMENT_CATEGORIES / CHART_CONFIG branch.
    PIPELINE_CONFIG['industry_id_key'] = industry_id

    if extract:
        from acc_deck_fs_pkg.api_extractor_v2 import extract_data

        print("\n" + "=" * 60)
        print("STAGE 0: DATA EXTRACTION")
        print("=" * 60)
        print(f"  Industry: {industry_id}")

        results = extract_data(prod_session, qa_session, industry_id=industry_id,
                               output_dir=str(_WRITABLE_DIR))
        df_forecast = results.get('prod')
        df_actuals  = results.get('qa')
        print(f"\n  Extraction complete: prod={'OK' if df_forecast is not None else 'FAILED'}, "
              f"qa={'OK' if df_actuals is not None else 'FAILED'}")

        if df_forecast is None or df_actuals is None:
            print("  Warning: Missing data source. Deck may be incomplete.")

        # Pass DataFrames directly — no CSV read required
        return main(df_forecast=df_forecast, df_actuals=df_actuals)

    return main()


if __name__ == "__main__":
    main()
