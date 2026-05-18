"""
main_meta_modes.py
==================
Main pipeline for Forecast Accuracy Deck Builder.

Generates PowerPoint presentations with LLM-powered insights comparing
forecast vs actual performance across categories.

Key Features
------------
- Direct meta-insight generation
- Configurable category ordering (sales volume/alphabetical)
- External prompt management via JSON config
- Sample insight injection for prompt engineering
- Grammar refinement for output quality (uses Opus for best compliance)
- Error recovery: skips failed categories and continues processing
- Cancellation support: checks for cancellation flag during execution

Usage
-----
This module is designed to be called from the GUI via pipeline_runner.py:

    from acc_deck_pkg.main_meta_modes import main
    main(runtime_config=cfg, category_order="sales_volume")

The runtime_config dict must contain paths, API keys, and prompt configurations.
See config_loader.py for the expected structure.

Author: Forecast Accuracy Team
"""

from __future__ import annotations

import os
import warnings
warnings.filterwarnings("ignore")

from datetime import datetime
from pathlib import Path
from typing import TYPE_CHECKING, Optional, Callable

if TYPE_CHECKING:
    import pandas as pd
    from pptx import Presentation


class CancellationError(Exception):
    """Raised when user cancels pipeline execution."""
    pass


# Import config/prompt loaders (lightweight)
from acc_deck_pkg.prompt_loader import (
    _load_prompts_from_config,
    _append_examples_if_available,
    _get_api_model_and_timeout,
    _get_model_for_task,
)

# NOTE: Heavy imports (pandas, pptx) moved inside main() for faster startup
# Individual module imports still at top since they're needed by main()
from acc_deck_pkg.data_io import load_data
from acc_deck_pkg.yoy_transformers import yoy
from acc_deck_pkg.analysis import (
    build_topline,
    get_categories_by_sales_volume,
    make_tot_table,
    merge_analysis_tables,
)
from acc_deck_pkg import ppt_builder
from acc_deck_pkg.llm_insights_free import (
    generate_llm_insights_remote,
    generate_meta_slide_insights,
    generate_meta_slide_insights_from_data,
    generate_total_slide_subheader,
)


def _get_llm_fns(cfg: dict):
    """Return the four LLM-callable functions used by the pipeline.

    Tuple order is fixed: (meta_from_data, total_subheader, remote_insights,
    meta_insights). cfg is accepted for forward-compatibility — when the
    receiving team swaps the provider via src/llm/, no callers need to change.
    """
    return (
        generate_meta_slide_insights_from_data,
        generate_total_slide_subheader,
        generate_llm_insights_remote,
        generate_meta_slide_insights,
    )
from acc_deck_pkg.llm_grammar_refine import (
    refine_insight_optional,
    refine_meta_df,
    refine_meta_df_with_validation,
    refine_category_insights_df,
)
from acc_deck_pkg.slide_insight_adder import (
    map_categories_from_placeholders,
    filter_slide_mapping,
    apply_meta_insights,
    apply_total_subheader_to_slide,
)

BASE_DIR = Path(__file__).resolve().parent


def _load_sampled_examples(cfg: dict, sample_type: str = "meta", max_samples: int = 5) -> list:
    """
    Load director-approved example insights from CSV for prompt injection.

    Parameters
    ----------
    cfg : dict
        Runtime config with paths
    sample_type : str
        'meta' or 'total' - which sample CSV to load
    max_samples : int
        Maximum number of examples to return

    Returns
    -------
    list of dict
        Sampled example insights
    """
    import pandas as pd
    import random

    base_dir = Path(cfg.get("_config_dir", str(Path.cwd())))

    # Determine which CSV to load
    prompts_cfg = cfg.get("prompts", {})
    if sample_type == "meta":
        csv_key = "sample_insights_meta_csv"
    else:
        csv_key = "sample_insights_total_csv"

    csv_rel = prompts_cfg.get(csv_key)
    if not csv_rel:
        return []

    csv_path = (base_dir / csv_rel).resolve()
    if not csv_path.exists():
        print(f"Warning: Sample CSV not found: {csv_path}")
        return []

    try:
        # Try multiple encodings - Excel CSVs often use cp1252 or latin-1
        df = None
        encodings_to_try = ['cp1252', 'latin-1', 'iso-8859-1', 'utf-8']

        for encoding in encodings_to_try:
            try:
                df = pd.read_csv(csv_path, encoding=encoding)
                if df is not None and not df.empty:
                    break
            except (UnicodeDecodeError, UnicodeError, Exception):
                df = None
                continue

        # Last resort: read with encoding_errors='replace' (pandas 1.3+)
        if df is None:
            try:
                df = pd.read_csv(csv_path, encoding='utf-8', encoding_errors='replace')
            except TypeError:
                # Older pandas without encoding_errors param
                import io
                with open(csv_path, 'r', encoding='latin-1', errors='replace') as f:
                    df = pd.read_csv(io.StringIO(f.read()))

        if df is None or df.empty:
            return []

        # Sample randomly (no seed = different every time)
        n = min(max_samples, len(df))
        sampled = df.sample(n=n)

        # Convert to list of dicts
        return sampled.to_dict('records')
    except Exception as e:
        print(f"Warning: Failed to load samples from {csv_path}: {e}")
        return []


def main(
        category_order: str = "sales_volume",
        runtime_config: Optional[dict] = None,
        cancel_check: Optional[Callable[[], bool]] = None,
        df=None,
) -> dict:
    """
    Main deck generation pipeline with error recovery.

    Generates a PowerPoint presentation with LLM-powered insights comparing
    forecast vs actual performance. Supports graceful error handling - if
    individual categories fail, processing continues with remaining categories.

    Parameters
    ----------
    category_order : str, default "sales_volume"
        How to order categories in the deck:
        - 'sales_volume': Largest sales first (recommended)
        - 'alphabetical': A-Z ordering
    runtime_config : dict, required
        Configuration dict from GUI containing:
        - paths: Input file paths (actual, forecast)
        - api_key: legacy slot (ignored — providers resolve their own keys
          via src/llm/providers/, with kwargs overrides for groq/moonshot)
        - prompts: Prompt file paths
        - deck_path: Output PowerPoint path
        - out_xlsx: Output Excel path
        - Plus column mappings and model parameters
    cancel_check : callable, optional
        Function that returns True if user requested cancellation.
        Called periodically during processing to allow early termination.

    Returns
    -------
    dict
        Execution summary containing:
        - success: bool - Overall success status
        - categories_processed: int - Number successfully processed
        - categories_failed: int - Number that failed (skipped)
        - failed_categories: list - Names of failed categories
        - warnings: list - Non-fatal warnings encountered

    Raises
    ------
    ValueError
        If runtime_config is None or missing required keys.
    CancellationError
        If cancel_check returns True during execution.
    """
    # LAZY IMPORTS: Load heavy libraries only when pipeline runs
    import pandas as pd
    from pptx import Presentation

    # Execution tracking
    execution_summary = {
        "success": False,
        "categories_processed": 0,
        "categories_failed": 0,
        "failed_categories": [],
        "warnings": [],
    }

    def check_cancelled():
        """Check if user requested cancellation."""
        if cancel_check and cancel_check():
            raise CancellationError("Pipeline cancelled by user")

    # Configure pandas display
    pd.set_option("display.max_columns", None)
    pd.set_option("display.width", None)
    pd.set_option("display.max_colwidth", None)

    # Validate config
    if runtime_config is None:
        raise ValueError(
            "runtime_config is required. This script is designed to be called "
            "from the GUI with a complete configuration."
        )

    cfg = runtime_config
    deck_path = cfg["deck_path"]
    out_xlsx = cfg["out_xlsx"]
    api_key = cfg["api_key"]

    # Resolve LLM functions. After the May 2026 cleanup there is only one
    # provider path (free-tier via src/llm/) — _get_llm_fns is kept as a
    # seam so a future swap to internally-hosted models doesn't ripple
    # through call sites.
    (
        _fn_meta_from_data,
        _fn_total_subheader,
        _fn_remote_insights,    # unused since traditional mode was dropped
        _fn_meta_insights,      # unused since traditional mode was dropped
    ) = _get_llm_fns(cfg)

    # Per-call API key overrides forwarded into every LLM function via
    # **kwargs. None / missing keys → providers fall back to env vars
    # (GROQ_API_KEY, MOONSHOT_API_KEY).
    _free_cfg = cfg.get("free_llm", {})
    _free_llm_kwargs = {
        "groq_api_key":     cfg.get("groq_api_key")     or os.getenv(_free_cfg.get("groq_api_key_env",     "GROQ_API_KEY"), "")     or None,
        "moonshot_api_key": cfg.get("moonshot_api_key") or os.getenv(_free_cfg.get("moonshot_api_key_env", "MOONSHOT_API_KEY"), "") or None,
    }

    # Load prompts from external files
    system_prompt, meta_prompt, total_prompt, row_prompt = _load_prompts_from_config(cfg)

    # Append sample insights (configurable)
    meta_prompt, total_prompt = _append_examples_if_available(
        cfg=cfg,
        meta_prompt=meta_prompt,
        total_prompt=total_prompt,
    )

    # Get default model and timeout from config
    default_model, timeout = _get_api_model_and_timeout(cfg)

    # Get per-task models (allows Opus for meta insights, Sonnet for others)
    model_category = _get_model_for_task(cfg, "category_insights")
    model_meta = _get_model_for_task(cfg, "meta_insights")
    model_total = _get_model_for_task(cfg, "total_slide")
    model_refine = _get_model_for_task(cfg, "grammar_refine")

    # Extract model parameters
    model_params = cfg.get("model_params", {})
    cat_params = model_params.get("category_insights", {})
    meta_params = model_params.get("meta_insights", {})
    total_params = model_params.get("total_slide", {})
    refine_params = model_params.get("grammar_refine", {})

    # Display configuration (condensed)
    _model_label = "GPT-OSS-120B (Groq) + Kimi K2.6 (Moonshot)"
    print(f"-- Pipeline: {category_order} | {_model_label} --")

    # === STAGE 1: SETUP ===
    print("\n[1/5] Setting up presentation...")
    prs = Presentation(cfg["ppt_template"])
    intro_slide = ppt_builder.add_slide(prs, "intro")
    type_map = {prs.slides.index(intro_slide): "intro"}

    # === STAGE 2: LOAD DATA ===
    print("[2/5] Loading data...")
    if df is None:
        df = load_data(cfg)
    else:
        print(f"      Using pre-loaded DataFrame ({len(df):,} rows)")
    print(f"      Loaded {len(df):,} rows")

    # === STAGE 3: BUILD CHARTS ===
    print("[3/5] Building charts...")

    # Industry-specific rules: exclude_from_total / chart_excluded_categories
    _industry_rules = cfg.get("industry_rules", {})
    _level1 = cfg.get("input_level1", "")
    _active_rules = next(
        (v for k, v in _industry_rules.items() if k.lower() == _level1.lower()),
        {}
    )
    _exclude_from_total = _active_rules.get("exclude_from_total", [])
    _chart_excluded     = _active_rules.get("chart_excluded_categories", [])

    if _exclude_from_total:
        print(f"      Industry rule: excluding {_exclude_from_total} from total")
    if _chart_excluded:
        print(f"      Industry rule: {_chart_excluded} included in total only (no chart)")

    df_for_total = df[~df["level2"].isin(_exclude_from_total)].copy() if _exclude_from_total else df
    df_tot = make_tot_table(df_for_total, cfg["input_year"], cfg["input_quarter"], cfg["input_level1"])
    print(df_tot)
    tot_slides = ppt_builder.add_slides(prs, df_tot, "tot")
    for s in tot_slides:
        type_map[prs.slides.index(s)] = "tot"
    tot_charts = ppt_builder.add_charts(prs, df_tot, tot_slides)

    # Get categories in requested order
    topline_calc = build_topline(df, cfg["input_year"], cfg["input_quarter"])

    if category_order == "alphabetical":
        ordered_categories = sorted(df["level2"].dropna().unique().tolist())
    else:
        ordered_categories = get_categories_by_sales_volume(topline_calc, ascending=False)

    # Remove chart-excluded categories and total-excluded categories from charting
    _no_chart = set(_chart_excluded) | set(_exclude_from_total)
    if _no_chart:
        ordered_categories = [c for c in ordered_categories if c not in _no_chart]

    print(f"      {len(ordered_categories)} categories ({category_order} order)")

    # Process categories
    check_cancelled()
    all_category_dfs = []
    processed_categories = []

    for i, cat in enumerate(ordered_categories):
        if i > 0 and i % 5 == 0:
            check_cancelled()

        try:
            t = yoy(
                df,
                input_year=cfg["input_year"],
                input_quarter=cfg["input_quarter"],
                cat=cat,
                total=False,
                topline=False,
            )
            t.columns.name = cat
            all_category_dfs.append(t)
            processed_categories.append(cat)
            execution_summary["categories_processed"] += 1

        except Exception as e:
            warning_msg = f"Skipped category '{cat}': {str(e)}"
            execution_summary["categories_failed"] += 1
            execution_summary["failed_categories"].append(cat)
            execution_summary["warnings"].append(warning_msg)
            continue

    if not all_category_dfs:
        raise ValueError(
            f"All {len(ordered_categories)} categories failed to process. "
            f"Check data format and column mappings."
        )

    if execution_summary["categories_failed"] > 0:
        print(f"      Warning: {execution_summary['categories_failed']} categories skipped")

    # Add category slides and charts
    category_slides = ppt_builder.add_slides(prs, all_category_dfs, "category")
    for s in category_slides:
        type_map[prs.slides.index(s)] = "category"
    category_charts = ppt_builder.add_charts(prs, all_category_dfs, category_slides)
    print(f"      Created {len(category_slides)} slides with charts")

    # Apply titles, footers, chart styling
    ppt_builder.apply_titles(prs, cfg, type_dict_pres=type_map, today=datetime.now())
    pres_charts = tot_charts + category_charts
    pres_dfs = df_tot + all_category_dfs
    ppt_builder.format_charts(prs, pres_charts, pres_dfs)

    # === STAGE 4: GENERATE INSIGHTS ===
    print("[4/5] Generating insights...")

    # Total slide insight
    total_text = None
    try:
        total_text = _fn_total_subheader(
            df_tot=df_tot,
            system_prompt=system_prompt,
            user_total_prompt=total_prompt,
            api_key=api_key,
            model=model_total,
            timeout=timeout,
            temperature=total_params.get("temperature", 0.65),
            top_p=total_params.get("top_p", 0.9),
            max_tokens=total_params.get("max_tokens", 250),
            **_free_llm_kwargs,
        )

        total_text = refine_insight_optional(
            total_text,
            kind="total",
            api_key=api_key,
            model=model_refine,
            timeout=timeout,
            temperature=refine_params.get("temperature", 0.5),
            top_p=refine_params.get("top_p", 0.82),
            max_tokens=refine_params.get("max_tokens", 180),
            force=True,
        )

        _ = apply_total_subheader_to_slide(
            prs=prs,
            total_slide=tot_slides[0],
            subheader_text=total_text,
            placeholder_idx=14,
        )
    except Exception as e:
        print(f"      Warning: Total insight skipped - {e}")

    # Prepare category data
    check_cancelled()
    collapsed = merge_analysis_tables(all_category_dfs)

    if collapsed.empty:
        prs.save(deck_path)
        print(f"      Saved: {deck_path}")
        execution_summary["success"] = True
        return execution_summary

    # Map slides to categories
    slide_cats = map_categories_from_placeholders(prs)
    filtered_slide_cats = filter_slide_mapping(slide_cats, collapsed)

    # Generate meta-insights
    check_cancelled()
    meta_df = None
    insights_df = collapsed

    try:
        sampling_cfg = cfg.get("prompt_sampling", {})
        max_samples = sampling_cfg.get("max_samples_meta", 5)
        sampled_examples = _load_sampled_examples(cfg, sample_type="meta", max_samples=max_samples)

        # Load narrative analysis config
        narrative_config = cfg.get("narrative_analysis", {})

        meta_df = _fn_meta_from_data(
            slide_mapping=filtered_slide_cats,
            collapsed_df=collapsed,
            system_prompt=system_prompt,
            user_meta_prompt=meta_prompt,
            api_key=api_key,
            model=model_meta,
            timeout=timeout,
            temperature=meta_params.get("temperature", 0.8),
            top_p=meta_params.get("top_p", 0.79),
            max_tokens=meta_params.get("max_tokens", 300),
            sampled_examples=sampled_examples,
            narrative_config=narrative_config,
            **_free_llm_kwargs,
        )
        insights_df = collapsed

    except CancellationError:
        raise  # Re-raise cancellation errors
    except Exception as e:
        warning_msg = f"Meta-insight generation failed: {str(e)}"
        print(f"Warning: {warning_msg}")
        execution_summary["warnings"].append(warning_msg)
        # Create empty meta_df so we can still save the deck
        import pandas as pd
        meta_df = pd.DataFrame({
            "slide_id": list(filtered_slide_cats.keys()),
            "meta_insight": ["(Insight generation failed)" for _ in filtered_slide_cats]
        })

    # Clean up meta-insights
    meta_df["meta_insight"] = meta_df["meta_insight"].astype(str).str.strip(' "\'""\'')

    # Refinement pass with separate style examples
    check_cancelled()

    # Load fresh samples for refinement (different from generation samples)
    sampling_cfg = cfg.get("prompt_sampling", {})
    max_refine_samples = sampling_cfg.get("max_samples_refine", 3)
    refine_examples = _load_sampled_examples(cfg, sample_type="meta", max_samples=max_refine_samples)

    try:
        meta_df = refine_meta_df_with_validation(
            meta_df,
            collapsed_df=collapsed,
            slide_mapping=filtered_slide_cats,
            sampled_examples=refine_examples,
            api_key=api_key,
            model=model_refine,
            timeout=timeout,
            temperature=refine_params.get("temperature", 0.3),
            top_p=refine_params.get("top_p", 0.92),
            max_tokens=refine_params.get("max_tokens", 180),
            verbose=False,
        )
    except CancellationError:
        raise
    except Exception as e:
        warning_msg = f"Refinement skipped: {str(e)}"
        execution_summary["warnings"].append(warning_msg)

    # Apply meta-insights to slides
    meta_status = apply_meta_insights(
        prs=prs,
        meta_df=meta_df,
        placeholder_idx=14,
    )

    # === STAGE 5: SAVE & SUMMARY ===
    print("\n[5/5] Saving outputs...")

    # Save workbook - organized by slide with meta insight and category data
    with pd.ExcelWriter(out_xlsx, engine="openpyxl") as writer:
        # Build slide-focused data: each row is a category, with meta insight on every row
        slide_data_rows = []

        for slide_id, categories in filtered_slide_cats.items():
            # Get the meta insight for this slide
            insight_row = meta_df[meta_df["slide_id"] == slide_id]
            meta_insight = insight_row["meta_insight"].values[0] if len(insight_row) > 0 else ""

            # Get category data for this slide
            for cat in categories:
                cat_data = collapsed[collapsed["category"] == cat]
                if cat_data.empty:
                    continue

                row = cat_data.iloc[0]
                slide_data_rows.append({
                    "slide": slide_id,
                    "category": cat,
                    "dollars_var": row.get("Diff (%)_Dollars", ""),
                    "units_var": row.get("Diff (%)_Units", ""),
                    "asp_var": row.get("Diff (%)_ASP", ""),
                    "slide_insight": meta_insight,  # Show on every row for this slide
                })

        slide_df = pd.DataFrame(slide_data_rows)
        slide_df.to_excel(writer, sheet_name="insights_by_slide", index=False)

    # Footer with trademark and quarter
    _ = ppt_builder.set_text_in_placeholder_31(
        prs,
        f"Future of\u2122 insights for {cfg['input_level1']} - Q{cfg['input_quarter']} {cfg['input_year']} Accuracy",
        mode="replace",
    )

    # Remove empty placeholders for cleaner client-ready output
    ppt_builder.remove_empty_placeholders(prs)

    # Save deck
    prs.save(deck_path)

    # === FINAL SUMMARY ===
    print("\n" + "=" * 60)
    print("GENERATED INSIGHTS SUMMARY")
    print("=" * 60)

    # Show total insight first
    if total_text:
        print(f"\nTOTAL SLIDE:")
        print(f"  \"{total_text}\"")

    # Show slide insights organized by slide
    print(f"\nSLIDE INSIGHTS ({len(meta_df)} slides):")
    for _, r in meta_df.iterrows():
        slide_id = r['slide_id']
        insight = str(r['meta_insight'])
        categories = filtered_slide_cats.get(slide_id, [])
        cat_str = ", ".join(categories[:3])
        if len(categories) > 3:
            cat_str += f" +{len(categories)-3} more"

        print(f"\n{slide_id}")
        print(f"  Categories: {cat_str}")
        print(f"  Insight: \"{insight}\"")

    print("\n" + "=" * 60)
    print(f"Saved: {Path(deck_path).name}")
    print(f"Excel: {Path(out_xlsx).name}")
    print("=" * 60)

    # Return execution summary
    execution_summary["success"] = True
    if execution_summary["warnings"]:
        print(f"\nCompleted with {len(execution_summary['warnings'])} warning(s)")
    return execution_summary


if __name__ == "__main__":
    print("Error: This script requires runtime_config from the GUI.")
    print("        Please run via run_gui.py")

